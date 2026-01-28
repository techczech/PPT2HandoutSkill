#!/usr/bin/env python3
"""
PPTX to JSON Extractor for Handout Site Generator

Extracts content from PowerPoint presentations into a JSON structure
compatible with the React handout site template.

Usage:
    python extract-pptx.py <input.pptx> [output_dir]
    python extract-pptx.py <input.pptx> [output_dir] --analyze-images

Output:
    output_dir/
    ├── presentation.json    (structured content)
    └── media/
        └── {uuid}/          (extracted images and videos)

Features:
    - Text extraction with formatting (bold, italic, underline, colors, hyperlinks)
    - Auto shape extraction (arrows, connectors, symbols like ≠, +, -, etc.)
    - Shape properties (position, colors, rotation, animation order)
    - Image extraction with optional AI-powered analysis
    - Table extraction
    - Speaker notes extraction
    - SmartArt diagram extraction (node hierarchy from diagram XML)
    - Section detection from slide layout names and content

Text Formatting:
    List items and headings include a 'runs' array when formatting is present:
    [{"text": "plain"}, {"text": "bold", "bold": true}, {"text": "text"}]

Shape Extraction:
    Meaningful shapes (arrows, symbols, connectors) are extracted with:
    - shape_type: auto shape type (e.g., "right_arrow", "math_not_equal")
    - shape_name: PowerPoint shape name
    - position: {left, top, width, height} in EMUs
    - fill_color, line_color: hex colors
    - rotation: angle in degrees
    - animation_order: entry order in animations (1-based)

Limitations:
    - Complex animations are simplified to just entry order
    - Image analysis recommended after extraction for descriptions
    - Manual review recommended after extraction
"""

import sys
import os
import json
import uuid
import re
import base64
import argparse
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.dml import MSO_THEME_COLOR
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

# Optional: Claude API for image analysis
ANTHROPIC_AVAILABLE = False
anthropic_client = None
try:
    import anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    pass

# Global flag for image analysis
ANALYZE_IMAGES = False


def init_anthropic():
    """Initialize Anthropic client if API key is available."""
    global anthropic_client
    if ANTHROPIC_AVAILABLE and os.environ.get('ANTHROPIC_API_KEY'):
        anthropic_client = anthropic.Anthropic()
        return True
    return False


def analyze_image(image_blob, ext, slide_title=None):
    """
    Analyze image using Claude's vision API.
    Returns a dict with description and optional quote extraction.
    """
    if not anthropic_client:
        return None

    # Map extensions to media types
    media_type_map = {
        'png': 'image/png',
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'gif': 'image/gif',
        'webp': 'image/webp'
    }

    media_type = media_type_map.get(ext.lower())
    if not media_type:
        return None

    try:
        # Encode image to base64
        image_data = base64.standard_b64encode(image_blob).decode('utf-8')

        context = f" from slide '{slide_title}'" if slide_title else ""

        response = anthropic_client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=500,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": image_data
                            }
                        },
                        {
                            "type": "text",
                            "text": f"""Analyze this presentation slide image{context}. Respond in JSON format:

{{
  "description": "Brief description of what the image shows",
  "has_quote": true/false,
  "quote_text": "The verbatim quote/message text if present",
  "quote_attribution": "Who said it (name, @handle, title)"
}}

Focus on:
1. Tweets - extract the full tweet text and @handle
2. Slack/chat messages - extract the message and sender name
3. Screenshots with quotes - extract any quotable statements
4. Social media posts - extract the post content and author

Examples:
- Tweet: {{"description": "Tweet screenshot", "has_quote": true, "quote_text": "ChatGPT launched on wednesday. today it crossed 1 million users!", "quote_attribution": "Sam Altman (@sama)"}}
- Slack: {{"description": "Slack message screenshot", "has_quote": true, "quote_text": "ChatGPT went viral. 100k people have tried this so far.", "quote_attribution": "Evan Morikawa, OpenAI"}}
- Chart: {{"description": "Bar chart showing AI adoption rates", "has_quote": false}}

Return ONLY valid JSON, no other text."""
                        }
                    ]
                }
            ]
        )

        # Parse JSON response
        response_text = response.content[0].text.strip()
        # Handle potential markdown code blocks
        if response_text.startswith('```'):
            response_text = response_text.split('```')[1]
            if response_text.startswith('json'):
                response_text = response_text[4:]
            response_text = response_text.strip()

        result = json.loads(response_text)
        return result

    except json.JSONDecodeError as e:
        print(f"    Warning: Could not parse image analysis JSON: {e}")
        # Return basic description if JSON parsing fails
        return {"description": response.content[0].text.strip() if response else None}
    except Exception as e:
        print(f"    Warning: Image analysis failed: {e}")
        return None


def slugify(text):
    """Convert text to a URL-friendly slug."""
    text = text.lower().strip()
    text = re.sub(r'[^\w\s-]', '', text)
    text = re.sub(r'[-\s]+', '-', text)
    return text[:50]


def extract_formatted_runs(paragraph):
    """Extract ALL text runs with formatting from a paragraph.

    Returns list of TextRun dicts only if there's any formatting,
    otherwise returns empty list (plain text is already in 'text' field).
    """
    runs = []
    has_any_formatting = False

    try:
        for run in paragraph.runs:
            if not run.text:
                continue

            # Get formatting
            bold = run.font.bold if run.font.bold is not None else False
            italic = run.font.italic if run.font.italic is not None else False
            underline = run.font.underline is not None and run.font.underline

            # Get URL if present
            url = None
            if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                url = run.hyperlink.address

            # Get font color
            font_color = None
            try:
                if run.font.color and run.font.color.rgb:
                    font_color = str(run.font.color.rgb)
            except Exception:
                pass

            # Track if any run has formatting
            if bold or italic or underline or url or font_color:
                has_any_formatting = True

            # Include ALL runs to preserve full text
            run_dict = {'text': run.text}
            if bold:
                run_dict['bold'] = True
            if italic:
                run_dict['italic'] = True
            if underline:
                run_dict['underline'] = True
            if url:
                run_dict['url'] = url
            if font_color:
                run_dict['font_color'] = font_color
            runs.append(run_dict)
    except Exception as e:
        print(f"    Warning: Could not extract formatted runs: {e}")

    return runs if has_any_formatting else []


def extract_animation_map(slide):
    """Extract animation order for shapes on a slide.

    Returns dict mapping shape_id -> animation_order (1-based).
    """
    animation_map = {}
    try:
        slide_elem = slide._element
        ns = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }

        timing = slide_elem.find('.//p:timing', ns)
        if timing is None:
            return animation_map

        seq = timing.find('.//p:seq', ns)
        if seq is None:
            return animation_map

        order = 0
        for child_tn_lst in seq.findall('.//p:childTnLst', ns):
            for par in child_tn_lst.findall('p:par', ns):
                for tgt_el in par.findall('.//p:tgtEl', ns):
                    sp_tgt = tgt_el.find('p:spTgt', ns)
                    if sp_tgt is not None:
                        shape_id_str = sp_tgt.get('spid')
                        if shape_id_str:
                            order += 1
                            shape_id = int(shape_id_str)
                            if shape_id not in animation_map:
                                animation_map[shape_id] = order
    except Exception as e:
        print(f"    Warning: Could not extract animation map: {e}")

    return animation_map


def extract_auto_shape(shape, animation_map):
    """Extract auto shape (arrow, connector, symbol, etc.).

    Returns ShapeBlock dict if this is a meaningful auto shape, None otherwise.
    """
    try:
        shape_type_val = shape.shape_type

        shape_type_names = {
            1: "auto_shape",
            2: "callout",
            5: "freeform",
            9: "line",
            19: "text_box",
            21: "connector",
        }

        # Skip types we handle elsewhere
        skip_types = {6, 7, 13, 14, 16, 17, 18, 19}
        if shape_type_val in skip_types:
            return None

        shape_name = shape.name if hasattr(shape, 'name') else ""

        meaningful_keywords = [
            'arrow', 'connector', 'line', 'equal', 'plus', 'minus',
            'chevron', 'block', 'star', 'heart', 'lightning', 'sun',
            'callout', 'bubble', 'cloud', 'oval', 'rectangle', 'triangle',
            'pentagon', 'hexagon', 'cross', 'notequal', 'not equal'
        ]

        name_lower = shape_name.lower()
        is_meaningful = any(kw in name_lower for kw in meaningful_keywords)

        auto_shape_type = None
        if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type:
            auto_shape_type = str(shape.auto_shape_type).split('.')[-1].lower()
            is_meaningful = is_meaningful or auto_shape_type not in ('rectangle', 'rounded_rectangle')

        if not is_meaningful and shape_type_val == 1:
            try:
                if hasattr(shape, 'fill') and shape.fill:
                    fill_type = shape.fill.type
                    if fill_type is None:
                        return None
            except Exception:
                return None

        if not is_meaningful:
            return None

        # Get position
        left = shape.left if hasattr(shape, 'left') and shape.left else 0
        top = shape.top if hasattr(shape, 'top') and shape.top else 0
        width = shape.width if hasattr(shape, 'width') and shape.width else 0
        height = shape.height if hasattr(shape, 'height') and shape.height else 0

        # Get text inside shape
        text = ""
        runs = []
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text.strip()
            for p in shape.text_frame.paragraphs:
                runs.extend(extract_formatted_runs(p))

        # Get colors
        fill_color = None
        line_color = None
        try:
            if hasattr(shape, 'fill') and shape.fill and shape.fill.fore_color:
                if shape.fill.fore_color.rgb:
                    fill_color = str(shape.fill.fore_color.rgb)
        except Exception:
            pass

        try:
            if hasattr(shape, 'line') and shape.line and shape.line.color:
                if shape.line.color.rgb:
                    line_color = str(shape.line.color.rgb)
        except Exception:
            pass

        # Get rotation
        rotation = 0.0
        if hasattr(shape, 'rotation'):
            rotation = shape.rotation or 0.0

        # Get animation order
        animation_order = None
        shape_id = shape.shape_id if hasattr(shape, 'shape_id') else None
        if shape_id and shape_id in animation_map:
            animation_order = animation_map[shape_id]

        result = {
            'type': 'shape',
            'shape_type': auto_shape_type or shape_type_names.get(shape_type_val, "shape"),
            'shape_name': shape_name,
            'position': {
                'left': left,
                'top': top,
                'width': width,
                'height': height,
            },
        }
        if text:
            result['text'] = text
        if runs:
            result['runs'] = runs
        if fill_color:
            result['fill_color'] = fill_color
        if line_color:
            result['line_color'] = line_color
        if rotation:
            result['rotation'] = rotation
        if animation_order is not None:
            result['animation_order'] = animation_order

        return result

    except Exception as e:
        print(f"    Warning: Could not extract auto shape: {e}")
        return None


def extract_smart_art(shape, slide, media_dir):
    """Extract SmartArt diagram content from a graphicFrame shape.

    Extracts node hierarchy, layout name, and icon images (sa_ prefix files).
    Icons are associated with their owning data nodes via presentation
    relationship tracing (presOf, presParOf connections).

    Returns a smart_art content block with nodes, or None if not SmartArt.
    """
    from lxml import etree
    import os

    el = shape._element
    ns_dgm = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    dgm_els = el.findall('.//{%s}relIds' % ns_dgm)
    if not dgm_els:
        return None

    try:
        # Get diagram data relationship
        dm_rid = dgm_els[0].get('{%s}dm' % ns_r)
        if not dm_rid or dm_rid not in slide.part.rels:
            return None

        rel = slide.part.rels[dm_rid]
        data_part = rel.target_part
        xml = etree.fromstring(data_part.blob)

        # Extract all diagram points (nodes) with icons
        pts = xml.findall('.//{%s}pt' % ns_dgm)
        cxns = xml.findall('.//{%s}cxn' % ns_dgm)

        # Build node map with icon extraction
        node_map = {}
        for pt in pts:
            mid = pt.get('modelId')
            ptype = pt.get('type', 'node')
            t_els = pt.findall('.//{%s}t' % ns_a)
            text = ' '.join((t.text or '') for t in t_els).strip()

            # Extract icon image from blip element
            icon = None
            icon_alt = None
            cnvpr = pt.find('.//{%s}cNvPr' % ns_a)
            if cnvpr is not None:
                icon_alt = cnvpr.get('descr') or cnvpr.get('title')
            blip = pt.find('.//{%s}blip' % ns_a)
            if blip is not None:
                rid = blip.get('{%s}embed' % ns_r)
                if rid:
                    try:
                        img_part = data_part.related_part(rid)
                        ext = img_part.content_type.split('/')[-1].replace('x-', '').replace('+xml', '')
                        safe_mid = mid.replace('{', '').replace('}', '').replace('-', '')
                        fname = f"sa_{safe_mid}.{ext}"
                        fpath = os.path.join(str(media_dir), fname)
                        with open(fpath, 'wb') as f:
                            f.write(img_part.blob)
                        icon = f"./{fname}"
                    except (KeyError, IOError) as e:
                        print(f"    Warning: Could not extract SmartArt icon: {e}")

            node_map[mid] = {
                'id': mid, 'type': ptype, 'text': text,
                'children_ids': [], 'icon': icon, 'icon_alt': icon_alt
            }

        # Build presentation relationship maps for icon reassignment
        visual_to_data = {}
        visual_parent = {}
        visual_children = {}
        data_root_id = None

        for mid, node in node_map.items():
            if node['type'] == 'doc':
                data_root_id = mid
                break

        for cxn in cxns:
            src_id = cxn.get('srcId', '')
            dst_id = cxn.get('destId', '')
            ctype = cxn.get('type', 'parOf')
            if ctype in ('parOf', ''):
                if src_id in node_map:
                    node_map[src_id]['children_ids'].append(dst_id)
            elif ctype == 'presOf':
                visual_to_data[dst_id] = src_id
            elif ctype == 'presParOf':
                visual_parent[dst_id] = src_id
                visual_children.setdefault(src_id, []).append(dst_id)

        # presAssocID explicit associations
        for pt in pts:
            mid = pt.get('modelId')
            pr_set = pt.find('{%s}prSet' % ns_dgm)
            if pr_set is not None:
                assoc_id = pr_set.get('presAssocID')
                if assoc_id:
                    visual_to_data[mid] = assoc_id

        # Reassign icons from visual nodes to their data owners
        def find_data_owner(vid):
            curr = vid
            visited = set()
            while curr and curr not in visited:
                visited.add(curr)
                if curr in visual_to_data:
                    did = visual_to_data[curr]
                    if did != data_root_id:
                        return did
                parent = visual_parent.get(curr)
                if parent:
                    for sib in visual_children.get(parent, []):
                        if sib != curr and sib in visual_to_data:
                            did = visual_to_data[sib]
                            if did != data_root_id:
                                return did
                curr = parent
            return None

        for mid, node in node_map.items():
            if node['icon']:
                owner_id = find_data_owner(mid)
                if owner_id and owner_id != mid and owner_id in node_map:
                    if not node_map[owner_id]['icon']:
                        node_map[owner_id]['icon'] = node['icon']
                        node_map[owner_id]['icon_alt'] = node['icon_alt']
                        node['icon'] = None
                        node['icon_alt'] = None

        # Find doc node for tree building
        doc_node = data_root_id

        def build_node(mid, level=0):
            node = node_map.get(mid)
            if not node or node['type'] != 'node':
                return None
            if not node['text'] and not node['icon']:
                return None
            result = {
                'id': mid[:8],
                'text': node['text'],
                'level': level,
                'children': [],
                'icon': node['icon'],
                'icon_alt': node['icon_alt']
            }
            for child_id in node.get('children_ids', []):
                child = build_node(child_id, level + 1)
                if child:
                    result['children'].append(child)
            return result

        # Build tree from doc node's children
        nodes = []
        if doc_node and doc_node in node_map:
            for child_id in node_map[doc_node]['children_ids']:
                node = build_node(child_id, 0)
                if node:
                    nodes.append(node)

        if not nodes:
            # Fallback: collect all nodes with text or icons
            for mid, node in node_map.items():
                if node['type'] == 'node' and (node['text'] or node['icon']):
                    nodes.append({
                        'id': mid[:8],
                        'text': node['text'],
                        'level': 0,
                        'children': [],
                        'icon': node['icon'],
                        'icon_alt': node['icon_alt']
                    })

        if not nodes:
            return None

        # Try to detect layout type from diagram XML
        layout_name = ''
        try:
            lo_rid = dgm_els[0].get('{%s}lo' % ns_r)
            if lo_rid and lo_rid in slide.part.rels:
                lo_rel = slide.part.rels[lo_rid]
                lo_xml = etree.fromstring(lo_rel.target_part.blob)
                # Layout name is in <dgm:title val="..."/> child element
                title_el = lo_xml.find('{%s}title' % ns_dgm)
                if title_el is not None:
                    layout_name = title_el.get('val', '')
        except Exception:
            pass

        icon_count = sum(1 for n in nodes if n['icon'])
        if icon_count:
            print(f"    SmartArt '{layout_name}': {len(nodes)} nodes, {icon_count} icons")

        return {
            'type': 'smart_art',
            'layout': layout_name,
            'nodes': nodes
        }

    except Exception as e:
        print(f"    Warning: Could not extract SmartArt: {e}")
        return None


def extract_text_from_shape(shape):
    """Extract text content from a shape with formatting."""
    if not shape.has_text_frame:
        return None

    paragraphs = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            # Detect list items by bullet or level
            level = para.level if para.level else 0
            item = {
                'text': text,
                'level': level,
                'is_bullet': level > 0 or (para.font and para.font.bold is False)
            }
            # Extract formatted runs
            runs = extract_formatted_runs(para)
            if runs:
                item['runs'] = runs
            paragraphs.append(item)

    return paragraphs if paragraphs else None


def extract_emf_embedded_image(emf_data: bytes):
    """
    Extract embedded JPEG/PNG image from EMF+ (Enhanced Metafile Plus) format.

    EMF+ files often contain raster images embedded in GDIC comment records.
    This is common when screenshots are pasted into PowerPoint.

    Returns:
        Tuple of (image_bytes, extension) if found, None otherwise
    """
    import struct

    # EMF+ uses comment records (type 70) to store GDI+ data
    # Look for GDIC records which often contain embedded images
    pos = 0
    while pos < len(emf_data) - 8:
        try:
            record_type, record_size = struct.unpack('<II', emf_data[pos:pos+8])
        except struct.error:
            break

        if record_type == 70:  # EMR_COMMENT (may contain EMF+ or GDIC data)
            comment_data = emf_data[pos+8:pos+record_size]
            if len(comment_data) > 8:
                # Check for GDIC identifier (contains embedded images)
                identifier = comment_data[4:8]
                if identifier == b'GDIC':
                    # Search for JPEG signature (FFD8FF)
                    jpg_sig = b'\xff\xd8\xff'
                    jpg_pos = comment_data.find(jpg_sig)
                    if jpg_pos >= 0:
                        # Extract JPEG - find EOI marker (0xFFD9)
                        jpg_data = comment_data[jpg_pos:]
                        eoi_pos = jpg_data.find(b'\xff\xd9')
                        if eoi_pos > 0:
                            jpg_data = jpg_data[:eoi_pos+2]
                            return (jpg_data, 'jpg')

                    # Search for PNG signature
                    png_sig = b'\x89PNG\r\n\x1a\n'
                    png_pos = comment_data.find(png_sig)
                    if png_pos >= 0:
                        # Extract PNG - find IEND chunk
                        png_data = comment_data[png_pos:]
                        iend_pos = png_data.find(b'IEND')
                        if iend_pos > 0:
                            png_data = png_data[:iend_pos+8]  # Include IEND + CRC
                            return (png_data, 'png')

        if record_type == 14 or record_size == 0:  # EOF or invalid
            break
        pos += record_size

    return None


def extract_image(shape, media_dir, slide_num, shape_idx, slide_title=None):
    """Extract image from shape and save to media directory."""
    try:
        if hasattr(shape, 'image') and hasattr(shape.image, 'blob'):
            image = shape.image
            ext = image.ext
            blob = image.blob
            content_type = getattr(image, 'content_type', '')

            # Handle EMF/WMF vector formats - try to extract embedded images or convert
            if ext in ('emf', 'wmf') or content_type in ('image/x-emf', 'image/x-wmf'):
                converted = False
                # First try to extract embedded images from EMF+ format
                try:
                    embedded = extract_emf_embedded_image(blob)
                    if embedded:
                        blob, ext = embedded
                        converted = True
                        print(f"    Extracted embedded image from EMF+ format")
                except Exception as emf_err:
                    pass  # Silent fail, try Pillow next

                # Fallback to Pillow conversion
                if not converted:
                    try:
                        from PIL import Image
                        import io
                        img = Image.open(io.BytesIO(blob))
                        png_buffer = io.BytesIO()
                        img.save(png_buffer, format='PNG')
                        blob = png_buffer.getvalue()
                        ext = 'png'
                        print(f"    Converted {content_type or 'EMF/WMF'} to PNG")
                    except Exception as conv_err:
                        print(f"    Warning: Could not convert {ext} to PNG: {conv_err}")

            filename = f"slide_{slide_num}_{shape_idx}.{ext}"
            filepath = media_dir / filename

            with open(filepath, 'wb') as f:
                f.write(blob)

            # Analyze image with Claude if enabled
            description = None
            quote_text = None
            quote_attribution = None

            if ANALYZE_IMAGES:
                print(f"    Analyzing image: {filename}...")
                analysis = analyze_image(image.blob, ext, slide_title)
                if analysis:
                    description = analysis.get('description')
                    if analysis.get('has_quote') and analysis.get('quote_text'):
                        quote_text = analysis.get('quote_text')
                        quote_attribution = analysis.get('quote_attribution')
                        print(f"    -> Quote: \"{quote_text[:60]}...\" - {quote_attribution}")
                    elif description:
                        print(f"    -> {description[:80]}...")

            result = {
                'type': 'image',
                'src': f"./{filename}",
                'alt': shape.name or f"Slide {slide_num} image",
                'caption': None,
                'description': description
            }

            # Add quote fields if present
            if quote_text:
                result['quote_text'] = quote_text
                result['quote_attribution'] = quote_attribution

            return result
    except Exception as e:
        print(f"  Warning: Could not extract image: {e}")

    return None


def extract_layout_background(slide_layout, media_dir, slide_num):
    """Extract the largest image from a slide layout (background image).

    Title/final slides often have a background image in the slide layout
    rather than on the slide itself. This extracts that image.
    """
    try:
        largest = None
        largest_size = 0
        for shape in slide_layout.shapes:
            if hasattr(shape, 'image') and hasattr(shape.image, 'blob'):
                blob_size = len(shape.image.blob)
                # Pick the largest image (background), skip small logos
                if blob_size > largest_size:
                    largest_size = blob_size
                    largest = shape

        if largest and largest_size > 100000:  # >100KB, likely a background
            image = largest.image
            ext = image.ext
            filename = f"layout_bg_{slide_num}.{ext}"
            filepath = os.path.join(media_dir, filename)
            with open(filepath, 'wb') as f:
                f.write(image.blob)
            print(f"  Extracted layout background: {filename} ({largest_size} bytes)")
            return f"./{filename}"
    except Exception as e:
        print(f"  Warning: Could not extract layout background: {e}")
    return None


def extract_video(shape, media_dir, slide_num):
    """Extract video from a shape (embedded or external URL).

    Checks for a:videoFile in nvPicPr, nvSpPr, and element tree.
    Handles external URLs (YouTube etc.) and embedded p14:media blobs.

    Returns a video/link content block, or None.
    """
    try:
        element = shape._element if hasattr(shape, '_element') else None
        if element is None:
            return None

        ns = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        }

        # Look for videoFile in nvPicPr/nvPr (picture shapes)
        videoFile = None
        nvPr = None
        if hasattr(element, 'nvPicPr') and element.nvPicPr is not None:
            nvPr = element.nvPicPr.nvPr
            videoFile = nvPr.find('.//a:videoFile', ns)

        # Also check nvSpPr for other shape types
        if videoFile is None and hasattr(element, 'nvSpPr') and element.nvSpPr is not None:
            nvPr = element.nvSpPr.nvPr
            videoFile = nvPr.find('.//a:videoFile', ns)

        # Fallback: whole element tree
        if videoFile is None:
            videoFile = element.find('.//a:videoFile', ns)
            if videoFile is not None:
                nvPr = element.find('.//p:nvPr', ns)

        if videoFile is None:
            return None

        video_title = shape.name if hasattr(shape, 'name') else "Video"

        # External video URL (YouTube etc.)
        video_link_rId = videoFile.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link')
        if video_link_rId:
            try:
                target = shape.part.target_ref(video_link_rId)
                if target and (target.startswith('http://') or target.startswith('https://')):
                    print(f"  Found external video: {target}")
                    return {
                        'type': 'video',
                        'src': target,
                        'title': video_title,
                        'external': True
                    }
            except Exception:
                pass

        # Embedded video (p14:media)
        if nvPr is not None:
            p14_media = nvPr.find('.//p14:media', ns)
            if p14_media is not None:
                embed_rId = p14_media.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if embed_rId:
                    try:
                        video_part = shape.part.related_part(embed_rId)
                        if video_part and hasattr(video_part, 'blob'):
                            content_type = video_part.content_type if hasattr(video_part, 'content_type') else 'video/mp4'
                            ext_map = {
                                'video/mp4': '.mp4',
                                'video/x-m4v': '.m4v',
                                'video/webm': '.webm',
                                'video/quicktime': '.mov',
                                'video/x-msvideo': '.avi',
                            }
                            ext = ext_map.get(content_type, '.mp4')
                            video_filename = f"slide_{slide_num}_{shape.shape_id}{ext}"
                            video_path = media_dir / video_filename
                            with open(video_path, 'wb') as f:
                                f.write(video_part.blob)
                            print(f"  Extracted embedded video: {video_filename}")
                            return {
                                'type': 'video',
                                'src': f"./{video_filename}",
                                'title': video_title
                            }
                    except Exception as e:
                        print(f"  Warning: Could not extract embedded video blob: {e}")

    except Exception as e:
        print(f"  Warning: Could not extract video info: {e}")
    return None


def process_slide(slide, slide_num, media_dir):
    """Process a single slide and extract its content."""
    content = []
    title = None
    notes = ""

    # Extract notes
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            notes = notes_slide.notes_text_frame.text.strip()

    # Extract animation map for this slide
    animation_map = extract_animation_map(slide)

    # Sort shapes by position (top, left) for consistent ordering
    shapes = list(slide.shapes)
    shapes.sort(key=lambda s: (
        s.top if (hasattr(s, 'top') and s.top is not None) else 0,
        s.left if (hasattr(s, 'left') and s.left is not None) else 0
    ))

    # Process shapes
    for idx, shape in enumerate(shapes):
        # Check for video in ANY shape type first (before other checks)
        video_content = extract_video(shape, media_dir, slide_num)
        if video_content:
            content.append(video_content)
            # Don't skip - shape may also have poster image or text

        # Title
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            if placeholder_type in [1, 3]:  # Title or Center Title
                if shape.has_text_frame:
                    title = shape.text_frame.text.strip()
                    # Clean up special characters
                    title = title.replace('\x0b', '\n').strip()
                continue  # Only skip title placeholders (1, 3)
            # Non-title placeholders (body, media, subtitle, etc.) fall through to content extraction

        # Images - use hasattr check to catch placeholder images (type 14) too
        if hasattr(shape, 'image') and hasattr(shape.image, 'blob'):
            img_content = extract_image(shape, media_dir, slide_num, idx, title)
            if img_content:
                content.append(img_content)
            continue

        # Tables (simplified extraction)
        if shape.has_table:
            table_text = []
            for row in shape.table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                table_text.append(' | '.join(row_text))
            if table_text:
                content.append({
                    'type': 'heading',
                    'text': '\n'.join(table_text),
                    'level': 3
                })
            continue

        # SmartArt diagrams (graphicFrame with diagram namespace)
        smart_art_block = extract_smart_art(shape, slide, media_dir)
        if smart_art_block:
            content.append(smart_art_block)
            continue

        # Auto shapes (arrows, connectors, symbols)
        shape_block = extract_auto_shape(shape, animation_map)
        if shape_block:
            content.append(shape_block)
            # If shape has text, also process as text (don't skip)
            if shape.has_text_frame and shape.text_frame.text.strip():
                text_content = extract_text_from_shape(shape)
                if text_content:
                    items = []
                    for p in text_content:
                        item = {'text': p['text'], 'children': []}
                        if 'runs' in p:
                            item['runs'] = p['runs']
                        items.append(item)
                    if items:
                        content.append({
                            'type': 'list',
                            'style': 'bullet',
                            'items': items
                        })
            continue

        # Text content
        text_content = extract_text_from_shape(shape)
        if text_content:
            # Check if it's a list or heading
            if len(text_content) == 1 and text_content[0]['level'] == 0:
                # Single paragraph - could be heading
                heading = {
                    'type': 'heading',
                    'text': text_content[0]['text'],
                    'level': 2
                }
                if 'runs' in text_content[0]:
                    heading['runs'] = text_content[0]['runs']
                content.append(heading)
            else:
                # Multiple items - treat as list
                items = []
                for p in text_content:
                    item = {'text': p['text'], 'children': []}
                    if 'runs' in p:
                        item['runs'] = p['runs']
                    items.append(item)
                if items:
                    content.append({
                        'type': 'list',
                        'style': 'bullet',
                        'items': items
                    })

    # Use actual PowerPoint layout name
    layout = slide.slide_layout.name

    result = {
        'order': slide_num,
        'title': title or f"Slide {slide_num}",
        'layout': layout,
        'notes': notes,
        'content': content
    }

    # Extract layout background image for title/final slides
    # These slides often have a background image in the slide layout, not the slide itself
    layout_lower = layout.lower()
    if 'title slide' in layout_lower or 'image background' in layout_lower:
        bg_image = extract_layout_background(slide.slide_layout, media_dir, slide_num)
        if bg_image:
            result['layout_background'] = bg_image

    return result


def extract_native_sections(prs):
    """Extract sections from the PPTX file's native section structure.

    PowerPoint stores sections in p14:sectionLst inside presentation.xml.
    Returns a list of {title, slide_ids} or None if no sections found.
    """
    from pptx.oxml.ns import qn

    # Build slide_id -> slide_order mapping
    slide_id_to_order = {}
    for idx, slide in enumerate(prs.slides, 1):
        slide_id_to_order[slide.slide_id] = idx

    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
    }

    try:
        ext_lst = prs.element.find('./p:extLst', namespaces=nsmap)
        if ext_lst is None:
            return None

        for ext in ext_lst.findall('./p:ext', namespaces=nsmap):
            section_list = ext.find('.//p14:sectionLst', namespaces=nsmap)
            if section_list is not None:
                sections = []
                for section_el in section_list.findall('./p14:section', namespaces=nsmap):
                    name = section_el.get('name', 'Untitled Section')
                    slide_orders = []
                    sld_id_lst = section_el.find('./p14:sldIdLst', namespaces=nsmap)
                    if sld_id_lst is not None:
                        for sld_id_tag in sld_id_lst.findall('./p14:sldId', namespaces=nsmap):
                            sid = int(sld_id_tag.get('id'))
                            if sid in slide_id_to_order:
                                slide_orders.append(slide_id_to_order[sid])
                    if slide_orders:
                        sections.append({'title': name, 'slide_orders': sorted(slide_orders)})
                if sections:
                    return sections
    except Exception as e:
        print(f"  Warning: Could not extract native sections: {e}")

    return None


def detect_sections(prs, slides_data):
    """Extract sections from PPTX file. Uses native PowerPoint sections if available,
    falls back to layout-based heuristic detection."""

    # Try native sections first
    native = extract_native_sections(prs)
    if native:
        print(f"  Using native PowerPoint sections ({len(native)} found)")
        # Build order -> slide_data lookup
        order_to_data = {s['order']: s for s in slides_data}

        sections = []
        for sec in native:
            section_slides = [order_to_data[o] for o in sec['slide_orders'] if o in order_to_data]
            if section_slides:
                sections.append({
                    'title': sec['title'],
                    'slides': section_slides
                })
        return sections

    # Fallback: heuristic detection from layout names
    print("  No native sections found, using layout-based detection")
    sections = []
    current_section = {
        'title': 'Main Content',
        'slides': []
    }

    for slide_data in slides_data:
        is_section_header = (
            'section heading' in slide_data['layout'].lower() or
            (len(slide_data['content']) == 0 and slide_data['title'] and
             any(kw in slide_data['layout'].lower() for kw in ['title only', 'title slide', 'full blue background']))
        )

        if is_section_header and current_section['slides']:
            sections.append(current_section)
            current_section = {
                'title': slide_data['title'],
                'slides': [slide_data]
            }
        else:
            current_section['slides'].append(slide_data)

    if current_section['slides']:
        sections.append(current_section)

    return sections


def extract_pptx(input_path, output_dir):
    """Main extraction function."""
    input_path = Path(input_path)
    output_dir = Path(output_dir)

    if not input_path.exists():
        print(f"Error: File not found: {input_path}")
        sys.exit(1)

    # Create output directories
    presentation_id = str(uuid.uuid4())
    media_dir = output_dir / "media" / presentation_id
    media_dir.mkdir(parents=True, exist_ok=True)

    print(f"Extracting: {input_path}")
    print(f"Output to: {output_dir}")
    print(f"Media dir: {media_dir}")
    print()

    # Load presentation
    prs = Presentation(input_path)

    # Process slides
    slides_data = []
    for idx, slide in enumerate(prs.slides, 1):
        print(f"Processing slide {idx}...")
        slide_data = process_slide(slide, idx, media_dir)
        slides_data.append(slide_data)
        if slide_data['title']:
            print(f"  Title: {slide_data['title'][:50]}...")

    print()

    # Detect sections
    sections = detect_sections(prs, slides_data)
    print(f"Detected {len(sections)} sections")

    # Count media files
    media_files = list(media_dir.glob('*'))
    image_count = len([f for f in media_files if f.suffix.lower() in ['.png', '.jpg', '.jpeg', '.gif']])
    video_count = len([f for f in media_files if f.suffix.lower() in ['.mp4', '.m4v', '.webm', '.mov', '.avi']])

    # Build output structure
    output = {
        'metadata': {
            'id': presentation_id,
            'source_file': input_path.name,
            'processed_at': datetime.now().isoformat(),
            'stats': {
                'slide_count': len(slides_data),
                'image_count': image_count,
                'video_count': video_count
            }
        },
        'sections': sections
    }

    # Write JSON
    json_path = output_dir / 'presentation.json'
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(output, f, indent=2, ensure_ascii=False)

    print()
    print("=" * 50)
    print("EXTRACTION COMPLETE")
    print("=" * 50)
    print(f"Slides extracted: {len(slides_data)}")
    print(f"Sections detected: {len(sections)}")
    print(f"Images extracted: {image_count}")
    print(f"Videos extracted: {video_count}")
    print(f"Output JSON: {json_path}")
    print(f"Media folder: {media_dir}")
    print()
    print("IMPORTANT: Please review the extracted content!")
    print("- Check slide titles are correct")
    print("- Verify section groupings make sense")
    print("- Review any images that may need manual adjustment")
    print("- SmartArt diagrams may need manual recreation")

    return output


def main():
    global ANALYZE_IMAGES

    parser = argparse.ArgumentParser(
        description=__doc__,
        formatter_class=argparse.RawDescriptionHelpFormatter
    )
    parser.add_argument('input', help='Input PowerPoint file (.pptx)')
    parser.add_argument('output', nargs='?', default='sourcematerials',
                        help='Output directory (default: sourcematerials)')
    parser.add_argument('--analyze-images', action='store_true',
                        help='Analyze images with Claude AI to generate descriptions. '
                             'Requires ANTHROPIC_API_KEY environment variable.')

    args = parser.parse_args()

    # Initialize image analysis if requested
    if args.analyze_images:
        if not ANTHROPIC_AVAILABLE:
            print("Error: --analyze-images requires the anthropic package.")
            print("Install with: pip install anthropic")
            sys.exit(1)

        if not os.environ.get('ANTHROPIC_API_KEY'):
            print("Error: --analyze-images requires ANTHROPIC_API_KEY environment variable.")
            sys.exit(1)

        if init_anthropic():
            ANALYZE_IMAGES = True
            print("Image analysis enabled (using Claude AI)")
            print()
        else:
            print("Warning: Could not initialize Anthropic client. Continuing without image analysis.")

    extract_pptx(args.input, args.output)


if __name__ == "__main__":
    main()
